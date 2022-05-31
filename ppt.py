from pptx import Presentation
from PIL import Image
import os
from pptx.slide import SlideLayout
from pptx.util import Inches, Pt, Cm
import datetime
import subprocess
import sys
import glob

#スライドサイズ
#4:3 (default) 9144000x6858000
#16:9 12193200x6858000
SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000
#スライド中心のX、Y座標（左上が原点）
IMG_CENTER_X, IMG_CENTER_Y = SLIDE_WIDTH/2, SLIDE_HEIGHT/2
#スライドのアスペクト比
SLIDE_ASPECT_RATIO = SLIDE_WIDTH / SLIDE_HEIGHT

#出力ファイル名
OUTPUT_FILE_PATH = "test.pptx"
#画像の格納ディレクトリ
IMG_DIR = "tmp"

def cmdret(result):
    ret = result.returncode
    if ret != 0:
      print(ret)
      print(result.stdout)
      print(result.stderr)
      sys.exit(ret)

def main():
    src = sys.argv[1]

    title = "テスト"
    target = "*.jpg"
    #rotate = "+90"
    #crop = "1936x1936+0+240"
    #resize = "784x784!"
    #resize = "1024x784!"
    #SLIDE_WIDTH, SLIDE_HEIGHT = 12193200, 6858000

    rotate = "+179"
    crop = "2950x2000+220+320"
    #resize = "1024x784!"
    #resize = "752x564!"
    #resize = "960x564!"

    now = datetime.datetime.now() # 現在時刻の取得
    today = now.strftime('%Y年%m月%d日') # 現在時刻を年月曜日で表示
    save_name = 'test.pptx' # 保存用のパワポのファイル名

    ## 作業ファイルをコピー
    result = subprocess.run(
        ["rm", "-fr", IMG_DIR],
        stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    cmdret(result)
    result = subprocess.run(
        ["cp", "-fr", src, IMG_DIR],
        stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    cmdret(result)

    ## ファイル一覧
    dir_path = IMG_DIR + "/" + target
    target_files = glob.glob(dir_path)
    #for file_name_i in target_files:
    #    print(file_name_i)

    ## ファイル一覧
    for file_name_i in target_files:
        #print(file_name_i)
        if 'rotate' in locals():
            result = subprocess.run(
                ["convert", file_name_i, "-rotate", rotate, file_name_i],
                stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            cmdret(result)
        result = subprocess.run(
            ["convert", file_name_i, "-crop", crop, file_name_i],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        cmdret(result)
        #result = subprocess.run(
        #    ["convert", file_name_i, "-resize", resize, file_name_i],
        #    stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
        #cmdret(result)


    #スライドオブジェクトの定義
    prs = Presentation()
    #スライドサイズの指定
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    #img画像のファイル名を取得
    img_files = os.listdir(IMG_DIR)
    #pngで終了するファイル名のみ抽出。貼り付けたい画像の拡張子に応じて変える
    img_files = [name for name in img_files if name.endswith(".jpg")]
    #昇順にソート（この順番でスライドに貼り付けられる）
    img_files.sort()#昇順にsort

    for name in img_files:
      print(name)
      path = IMG_DIR + "/" + name
      print(path)
      slide = add_slide(prs)
      add_picture(slide, path)
    #pptxファイルを出力する
    print(OUTPUT_FILE_PATH)
    prs.save(OUTPUT_FILE_PATH)


#受け取ったプレゼンテーションオブジェクトにスライドを追加し、追加されたスライドオブジェクトを返す。
def add_slide(prs):
  #白紙スライドの追加(ID=6は白紙スライド)
  blank_slide_layout = prs.slide_layouts[6]
  slide = prs.slides.add_slide(blank_slide_layout)
  return slide

#画像をスライド中心に貼り付ける
def add_picture(slide, img_file):
  #画像サイズを取得してアスペクト比を得る
  im = Image.open(img_file)
  im_width, im_height = im.size
  aspect_ratio = im_width/im_height

  #スライドと画像のアスペクト比に応じて処理を分岐
  #画像のほうが横長だったら横めいっぱいに広げる
  if aspect_ratio > SLIDE_ASPECT_RATIO:
    img_display_width = SLIDE_WIDTH
    img_display_height = img_display_width / aspect_ratio
  else: #画像のほうが縦長だったら縦めいっぱいに広げる
    img_display_height = SLIDE_HEIGHT
    img_display_width = img_display_height * aspect_ratio
  #センタリングする場合の画像の左上座標を計算
  left = IMG_CENTER_X - img_display_width / 2
  top = IMG_CENTER_Y - img_display_height / 2

  #画像をスライドに追加
  if aspect_ratio > SLIDE_ASPECT_RATIO:
    slide.shapes.add_picture(img_file, left, top, width = img_display_width)
  else:
    slide.shapes.add_picture(img_file, left, top, height = img_display_height)

  return slide

if __name__ == "__main__":
    main()

